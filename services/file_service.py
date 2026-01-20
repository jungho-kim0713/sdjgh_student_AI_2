"""
파일 업로드/텍스트 추출 관련 공용 유틸리티.
라우트에서 재사용하도록 함수 단위로 제공한다.
"""

from io import BytesIO

import pypdf
from docx import Document
from pptx import Presentation
import openpyxl

# 업로드 허용 확장자 목록
ALLOWED_EXTENSIONS = {
    "txt", "py", "md", "csv", "json", "html", "css", "js",
    "c", "cpp", "java", "pdf", "doc", "docx", "xls", "xlsx",
    "ppt", "pptx", "hwp", "png", "jpg", "jpeg", "gif", "webp"
}


def allowed_file(filename):
    """파일명 확장자가 허용 목록에 포함되는지 검사한다."""
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


def extract_text_from_file(file_content, filename):
    """
    파일 바이트에서 텍스트를 추출한다.
    이미지가 아닌 문서/오피스 파일을 대상으로 하며 실패 시 오류 메시지를 반환한다.
    """
    ext = filename.rsplit(".", 1)[1].lower()
    text = ""
    try:
        file_stream = BytesIO(file_content)
        if ext == "pdf":
            # PDF: 페이지별 텍스트 추출
            reader = pypdf.PdfReader(file_stream)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        elif ext in ["docx", "doc"]:
            # Word: 문단 단위 텍스트 추출
            doc = Document(file_stream)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif ext in ["pptx", "ppt"]:
            # PowerPoint: 슬라이드/도형 텍스트 추출
            prs = Presentation(file_stream)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext in ["xlsx", "xls"]:
            # Excel: 시트/행 단위 텍스트 추출
            wb = openpyxl.load_workbook(file_stream, data_only=True)
            for sheet in wb.worksheets:
                text += f"Sheet: {sheet.title}\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) for cell in row if cell is not None]
                    if row_text:
                        text += "\t".join(row_text) + "\n"
        else:
            # 기타 텍스트 파일: UTF-8로 디코딩
            text = file_content.decode("utf-8", errors="ignore")
        return text if text.strip() else "(내용 없음)"
    except Exception as e:
        return f"(텍스트 추출 실패: {e})"
